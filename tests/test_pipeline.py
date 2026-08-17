"""不依赖外部模型 API 的本地管线测试。"""

from __future__ import annotations

import hashlib

import numpy as np
from fastapi.testclient import TestClient
from langchain_core.documents import Document
from langchain_core.embeddings import Embeddings
from langchain_core.language_models import BaseChatModel
from langchain_core.messages import AIMessage
from langchain_core.outputs import ChatGeneration, ChatResult

from kb_agent import server
from kb_agent.chunking import split_documents
from kb_agent.ingest import ingest, load_documents
from kb_agent.loaders import load_file
from kb_agent.pdf_export import build_conversation_pdf, build_qa_pdf
from kb_agent.qa import (
    REPORT_SYSTEM_PROMPT,
    SYSTEM_PROMPT,
    UNKNOWN_ANSWER_WITH_OFFER,
    WEB_DISCLAIMER,
    WEB_SEARCH_FAILED,
    UNKNOWN_ANSWER,
    AnswerResult,
    GroundedRetriever,
    _conversation_queries,
    _complete_source_context,
    _extract_topic,
    _format_sources,
    _local_source_names,
    _is_other_region_source,
    _is_report_question,
    _is_summary_question,
    _needs_other_region_docs,
    _resolve_referents,
    _sanitize_answer,
    _should_web_search,
    answer_question,
    build_retriever,
    detect_target_sources,
    detect_tier,
    load_vectorstore,
)
from kb_agent.web_search import WebResult, search_web

# 测试夹具文件名统一放这里，避免在用例里散落写死；新增格式或改名时只改这一处。
SAMPLE_TXT = "sample.txt"
ENGLISH_RULES_TXT = "rules.txt"
FRUIT_TXT = "fruit.txt"
RULES_TXT = "规则.txt"
RULES_DOCX = "规则.docx"
RULES_DOC = "规则.doc"
RULES_PPTX = "规则.pptx"
RULES_PPT = "规则.ppt"
RULES_WPS = "规则.wps"
RULES_XLSX = "规则.xlsx"
RULES_CSV = "规则.csv"
RULES_PDF = "规则.pdf"
RULES_PNG = "规则.png"
SCANNED_PDF = "扫描件.pdf"
FEEDBACK_DOC = "大数据局反馈.doc"
FEEDBACK_0116_DOC = "大数据局反馈0116.doc"
PLAN_PDF = "计划.pdf"
MEASURE_DOCX = "办法.docx"
NATIONAL_TXT = "国家.txt"
PROVINCIAL_DOCX = "省级.docx"
MUNICIPAL_DOCX = "市级.docx"
REAL_TXT = "real.txt"
TEMP_DOCX = "~$临时.docx"
DEMO_TEST_PASSWORD = "2444666668888888"
AUTH_HEADERS = {"X-Demo-Password": DEMO_TEST_PASSWORD}


def _pdf_text(pdf_bytes: bytes) -> str:
    import pymupdf

    doc = pymupdf.open(stream=pdf_bytes, filetype="pdf")
    try:
        return "\n".join(page.get_text() for page in doc)
    finally:
        doc.close()


class HashEmbeddings(Embeddings):
    """确定性哈希向量，仅用于本地测试。"""

    def _embed(self, text: str) -> list[float]:
        digest = hashlib.sha256(text.encode("utf-8")).digest()
        vector = np.frombuffer(digest, dtype=np.uint8)[:16].astype(np.float64)
        return (vector / np.linalg.norm(vector)).tolist()

    def embed_documents(self, texts: list[str]) -> list[list[float]]:
        return [self._embed(text) for text in texts]

    def embed_query(self, text: str) -> list[float]:
        return self._embed(text)


class StubChatModel(BaseChatModel):
    """固定返回预设答案的假模型。"""

    answer: str = "3 天"

    @property
    def _llm_type(self) -> str:
        return "stub-chat-model"

    def _generate(self, messages, stop=None, run_manager=None, **kwargs) -> ChatResult:
        return ChatResult(generations=[ChatGeneration(message=AIMessage(content=self.answer))])


class RecordingChatModel(StubChatModel):
    """记录最后一次收到的消息，便于断言实际使用的提示词。"""

    last_messages: list = []

    def _generate(self, messages, stop=None, run_manager=None, **kwargs) -> ChatResult:
        self.last_messages = list(messages)
        return super()._generate(messages, stop=stop, run_manager=run_manager, **kwargs)


class StaticRetriever:
    """按问题关键字返回文档或空列表的假检索器。"""

    def __init__(self, documents: list[Document], unknown_questions: tuple[str, ...] = ()) -> None:
        self.documents = documents
        self.unknown_markers = unknown_questions

    def invoke(self, question: str) -> list[Document]:
        if any(marker in question for marker in self.unknown_markers):
            return []
        return self.documents


def test_ingest_and_retrieve(tmp_path) -> None:
    source_dir = tmp_path / "pdfs"
    source_dir.mkdir()
    (source_dir / SAMPLE_TXT).write_text("退货时间为 3 天。\n", encoding="utf-8")
    index_dir = tmp_path / "index"

    added = ingest(
        directory=source_dir,
        index_directory=index_dir,
        embedding=HashEmbeddings(),
        rebuild=True,
    )
    assert added == 1
    assert (index_dir / "index.faiss").exists()

    vectorstore = load_vectorstore(index_directory=index_dir, embedding=HashEmbeddings())
    retriever = build_retriever(vectorstore, top_k=1, score_threshold=0.0)
    documents = retriever.invoke("退货时间是几天？")
    assert documents
    assert "退货时间" in documents[0].page_content


def test_known_question_uses_retrieved_context() -> None:
    documents = [
        Document(
            page_content="退货时间为 3 天。",
            metadata={"source": RULES_TXT, "page": 1},
        )
    ]
    result = answer_question(
        "退货时间是几天？",
        retriever=StaticRetriever(documents),
        llm=StubChatModel(),
    )
    assert result.answer == "3 天"
    assert result.sources == [RULES_TXT]
    assert result.unknown is False


def test_unknown_question_returns_unknown() -> None:
    documents = [
        Document(
            page_content="退货时间为 3 天。",
            metadata={"source": RULES_TXT, "page": 1},
        )
    ]
    result = answer_question(
        "航母的航速是多少？",
        retriever=StaticRetriever(documents, unknown_questions=("航母",)),
        llm=StubChatModel(),
    )
    assert result.answer == UNKNOWN_ANSWER_WITH_OFFER
    assert result.sources == []
    assert result.unknown is True


def test_unknown_answer_hides_sources_even_with_retrieved_documents() -> None:
    documents = [
        Document(
            page_content="退货时间为 3 天。",
            metadata={"source": RULES_TXT, "page": 1},
        )
    ]
    result = answer_question(
        "退货时间是几天？",
        retriever=StaticRetriever(documents),
        llm=StubChatModel(answer=UNKNOWN_ANSWER),
    )
    assert result.answer == UNKNOWN_ANSWER_WITH_OFFER
    assert result.sources == []


def test_answer_question_includes_history_in_prompt() -> None:
    documents = [
        Document(
            page_content="退货时间为 3 天。",
            metadata={"source": RULES_TXT, "page": 1},
        )
    ]
    model = RecordingChatModel(answer="3 天")
    history = [
        {"role": "user", "content": "退货时间是几天？"},
        {"role": "assistant", "content": "退货时间为 3 天。"},
    ]
    result = answer_question(
        "那退款期限呢？",
        retriever=StaticRetriever(documents),
        llm=model,
        history=history,
    )
    assert result.answer == "3 天"
    assert result.unknown is False
    joined = "".join(str(message.content) for message in model.last_messages)
    assert "退货时间是几天？" in joined
    assert "退货时间为 3 天。" in joined


def test_answer_question_followup_uses_history_when_no_new_hits() -> None:
    history = [
        {"role": "user", "content": "苏州市数据条例对公共数据开放有什么规定？"},
        {
            "role": "assistant",
            "content": "公共数据开放遵循合法、公正、安全原则。",
            "sources": ["苏州市数据条例.docx"],
        },
    ]
    result = answer_question(
        "那它的适用范围呢？",
        retriever=StaticRetriever([]),
        llm=StubChatModel(answer="适用于行政机关和法律、法规授权的组织。"),
        history=history,
    )
    assert result.answer == "适用于行政机关和法律、法规授权的组织。"
    assert result.sources == ["苏州市数据条例.docx"]
    assert result.unknown is False


def test_answer_question_history_unknown_hides_sources() -> None:
    history = [
        {"role": "user", "content": "退货时间是几天？"},
        {
            "role": "assistant",
            "content": "退货时间为 3 天。",
            "sources": [RULES_TXT],
        },
    ]
    result = answer_question(
        "那退款金额是多少？",
        retriever=StaticRetriever([]),
        llm=StubChatModel(answer=UNKNOWN_ANSWER),
        history=history,
    )
    assert result.answer == UNKNOWN_ANSWER_WITH_OFFER
    assert result.sources == []
    assert result.unknown is True


def test_unknown_answer_includes_web_offer() -> None:
    result = answer_question(
        "航母的航速是多少？",
        retriever=StaticRetriever([]),
        llm=StubChatModel(),
    )
    assert result.answer == UNKNOWN_ANSWER_WITH_OFFER
    assert result.unknown is True
    assert result.web is False


def test_web_search_consent_after_unknown_uses_last_user_question(monkeypatch) -> None:
    captured: dict[str, object] = {}

    def fake_search(query: str, max_results: int = 5) -> list[WebResult]:
        captured["query"] = query
        return [
            WebResult(
                title="数字孪生政策摘要",
                url="https://example.com/digital-twin",
                snippet="公开网页中的数字孪生政策内容。",
            )
        ]

    monkeypatch.setattr("kb_agent.qa.search_web", fake_search)
    history = [
        {"role": "user", "content": "苏州市数字孪生政策是什么？"},
        {
            "role": "assistant",
            "content": UNKNOWN_ANSWER_WITH_OFFER,
            "sources": [],
            "unknown": True,
        },
    ]
    result = answer_question(
        "可以",
        retriever=StaticRetriever([]),
        llm=StubChatModel(answer="数字孪生相关政策摘要。"),
        history=history,
    )
    assert captured["query"] == "苏州市数字孪生政策是什么？"
    assert result.unknown is False
    assert result.web is True
    assert result.sources == []
    assert result.answer.startswith(WEB_DISCLAIMER)


def test_web_search_explicit_request_uses_question(monkeypatch) -> None:
    captured: dict[str, object] = {}

    def fake_search(query: str, max_results: int = 5) -> list[WebResult]:
        captured["query"] = query
        return [WebResult(title="标题", url="https://example.com", snippet="摘要")]

    monkeypatch.setattr("kb_agent.qa.search_web", fake_search)
    result = answer_question(
        "帮我联网搜索一下苏州市数据条例有哪些内容",
        retriever=StaticRetriever([]),
        llm=StubChatModel(answer="联网整理的内容。"),
    )
    assert captured["query"] == "苏州市数据条例有哪些内容"
    assert result.web is True
    assert result.unknown is False


def test_web_search_no_results_returns_failed_unknown(monkeypatch) -> None:
    monkeypatch.setattr(
        "kb_agent.qa.search_web",
        lambda query, max_results=5: [],
    )
    history = [
        {"role": "user", "content": "某个未知问题"},
        {"role": "assistant", "content": UNKNOWN_ANSWER_WITH_OFFER, "unknown": True},
    ]
    result = answer_question(
        "可以",
        retriever=StaticRetriever([]),
        llm=StubChatModel(),
        history=history,
    )
    assert result.answer == WEB_SEARCH_FAILED
    assert result.unknown is True
    assert result.web is True


def test_web_search_consent_requires_previous_unknown() -> None:
    assert _should_web_search("可以", []) is False
    assert (
        _should_web_search(
            "可以",
            [
                {"role": "user", "content": "问题"},
                {"role": "assistant", "content": "正常回答"},
            ],
        )
        is False
    )
    assert (
        _should_web_search(
            "可以",
            [
                {"role": "user", "content": "问题"},
                {"role": "assistant", "content": UNKNOWN_ANSWER_WITH_OFFER},
            ],
        )
        is True
    )
    assert _should_web_search("帮我搜索一下问题", []) is True


def test_web_search_parses_bing_rss(monkeypatch) -> None:
    xml = (
        '<?xml version="1.0"?><rss><channel><item><title>标题 &amp; 测试</title>'
        '<link>https://example.com/a</link>'
        "<description>&lt;p&gt;摘要内容&lt;/p&gt;</description></item></channel></rss>"
    ).encode("utf-8")

    class FakeResponse:
        content = xml
        text = ""

    monkeypatch.setattr("kb_agent.web_search._fetch", lambda url: FakeResponse())
    results = search_web("测试")
    assert results
    assert results[0].title == "标题 & 测试"
    assert results[0].url == "https://example.com/a"
    assert "摘要内容" in results[0].snippet


def test_extract_topic_gets_subject() -> None:
    assert _extract_topic("数字孪生是什么？") == "数字孪生"
    assert _extract_topic("苏州市数据条例对公共数据开放有什么规定？") == "苏州市数据条例"


def test_resolve_referents_replaces_pronoun_with_previous_topic() -> None:
    history = [
        {"role": "user", "content": "数字孪生是什么？"},
        {"role": "assistant", "content": "数字孪生是物理对象的数字化映射。"},
    ]
    assert _resolve_referents("园区对它有什么期望吗？", history) == "园区对数字孪生有什么期望吗？"
    assert _resolve_referents("那它呢？", history) == "那数字孪生呢？"
    assert _resolve_referents("该政策适用范围是什么？", history) == "数字孪生适用范围是什么？"


def test_conversation_queries_include_resolved_referent() -> None:
    history = [
        {"role": "user", "content": "数字孪生是什么？"},
        {"role": "assistant", "content": "数字孪生是物理对象的数字化映射。"},
    ]
    model = StubChatModel(answer="数字孪生 园区 期望")
    queries = _conversation_queries("园区对它有什么期望吗？", history, model)
    assert queries[0] == "园区对数字孪生有什么期望吗？"
    assert "它" not in queries[0]
    assert any("数字孪生" in query for query in queries)


def test_resolve_referents_with_multiple_questions_in_one_message() -> None:
    resolved = _resolve_referents("数字孪生是什么？园区对它有什么期望吗？", [])
    assert resolved == "数字孪生是什么？园区对数字孪生有什么期望吗？"


def test_resolve_referents_skips_referential_history_turn() -> None:
    history = [
        {"role": "user", "content": "数字孪生是什么？"},
        {"role": "assistant", "content": "数字孪生是物理对象的数字化映射。"},
        {"role": "user", "content": "园区对它有什么期望吗？"},
        {"role": "assistant", "content": "园区期望用于智慧治理。"},
    ]
    assert _resolve_referents("那它对产业有什么作用？", history) == "那数字孪生对产业有什么作用？"


def test_sanitize_answer_removes_knowledge_base_wording() -> None:
    assert _sanitize_answer("根据知识库内容，该政策适用园区。") == "根据我的理解，该政策适用园区。"
    assert _sanitize_answer("依据知识库中的文件") == "根据已有资料中的文件"
    assert "知识库" not in _sanitize_answer("知识库中没有相关规定。")
    assert _sanitize_answer("知识库中的《知识库管理办法》") == "已有资料中的《知识库管理办法》"


def test_prompts_forbid_knowledge_base_wording_in_answers() -> None:
    assert "不要出现“知识库”" in SYSTEM_PROMPT
    assert "根据我的理解" in SYSTEM_PROMPT
    assert "不要出现“知识库”" in REPORT_SYSTEM_PROMPT
    assert "根据我的理解" in REPORT_SYSTEM_PROMPT


def test_score_threshold_blocks_unrelated_question(tmp_path) -> None:
    source_dir = tmp_path / "pdfs"
    source_dir.mkdir()
    (source_dir / ENGLISH_RULES_TXT).write_text("退货时间为 3 天。\n", encoding="utf-8")
    (source_dir / FRUIT_TXT).write_text("苹果是红色的水果。\n", encoding="utf-8")
    index_dir = tmp_path / "index"
    ingest(
        directory=source_dir,
        index_directory=index_dir,
        embedding=HashEmbeddings(),
        rebuild=True,
    )
    vectorstore = load_vectorstore(index_directory=index_dir, embedding=HashEmbeddings())

    known = answer_question(
        "苹果是什么颜色？",
        vectorstore=vectorstore,
        llm=StubChatModel(answer="红色"),
        score_threshold=0.0,
    )
    assert known.sources
    assert known.unknown is False

    unknown = answer_question(
        "航母的航速是多少？",
        vectorstore=vectorstore,
        llm=StubChatModel(),
        score_threshold=1.01,
    )
    assert unknown.answer == UNKNOWN_ANSWER_WITH_OFFER
    assert unknown.sources == []
    assert unknown.unknown is True


def test_server_health_and_ask(monkeypatch) -> None:
    captured: dict[str, object] = {}

    def fake_answer(question, **kwargs):
        captured["question"] = question
        captured["history"] = kwargs.get("history")
        return AnswerResult(answer="不知道", sources=[])

    monkeypatch.setattr(
        server,
        "answer_question",
        fake_answer,
    )
    client = TestClient(server.app)

    assert client.get("/health").json() == {"status": "ok"}
    response = client.post(
        "/ask",
        headers=AUTH_HEADERS,
        json={
            "question": "未知问题",
            "history": [
                {"role": "user", "content": "退货时间是几天？"},
                {"role": "assistant", "content": "3 天"},
            ],
        },
    )
    assert response.status_code == 200
    assert response.json()["answer"] == "不知道"
    assert response.json()["unknown"] is False
    assert captured["question"] == "未知问题"
    assert captured["history"] == [
        {"role": "user", "content": "退货时间是几天？", "sources": [], "unknown": False},
        {"role": "assistant", "content": "3 天", "sources": [], "unknown": False},
    ]


def test_server_ask_returns_unknown_flag(monkeypatch) -> None:
    monkeypatch.setattr(
        server,
        "answer_question",
        lambda question, **kwargs: AnswerResult(
            answer=UNKNOWN_ANSWER_WITH_OFFER, sources=[], unknown=True
        ),
    )
    client = TestClient(server.app)
    response = client.post(
        "/ask", json={"question": "航母航速多少？"}, headers=AUTH_HEADERS
    )
    assert response.status_code == 200
    assert response.json()["answer"] == UNKNOWN_ANSWER_WITH_OFFER
    assert response.json()["sources"] == []
    assert response.json()["unknown"] is True


def test_server_ask_returns_web_flag(monkeypatch) -> None:
    monkeypatch.setattr(
        server,
        "answer_question",
        lambda question, **kwargs: AnswerResult(
            answer="联网答案", sources=[], unknown=False, web=True
        ),
    )
    client = TestClient(server.app)
    response = client.post("/ask", json={"question": "可以"}, headers=AUTH_HEADERS)
    assert response.status_code == 200
    assert response.json()["answer"] == "联网答案"
    assert response.json()["unknown"] is False
    assert response.json()["web"] is True


def test_server_ask_requires_password_when_enabled(monkeypatch) -> None:
    monkeypatch.setattr(server, "demo_password", lambda: DEMO_TEST_PASSWORD)
    monkeypatch.setattr(
        server,
        "answer_question",
        lambda question, **kwargs: AnswerResult(answer="ok", sources=[]),
    )
    client = TestClient(server.app)

    assert client.post("/ask", json={"question": "测试"}).status_code == 401
    assert (
        client.post(
            "/ask", json={"question": "测试"}, headers={"X-Demo-Password": "wrong"}
        ).status_code
        == 401
    )
    ok = client.post("/ask", json={"question": "测试"}, headers=AUTH_HEADERS)
    assert ok.status_code == 200
    assert ok.json()["answer"] == "ok"


def test_server_auth_status_and_check(monkeypatch) -> None:
    monkeypatch.setattr(server, "demo_password", lambda: DEMO_TEST_PASSWORD)
    client = TestClient(server.app)
    status = client.get("/auth/status").json()
    assert status["enabled"] is True
    assert status["locked"] is False
    assert (
        client.post(
            "/auth/check", json={"password": DEMO_TEST_PASSWORD}
        ).json()
        == {"ok": True}
    )
    assert client.post("/auth/check", json={"password": "wrong"}).status_code == 401

    monkeypatch.setattr(server, "demo_password", lambda: "")
    disabled = client.get("/auth/status").json()
    assert disabled["enabled"] is False
    assert disabled["locked"] is False
    monkeypatch.setattr(
        server,
        "answer_question",
        lambda question, **kwargs: AnswerResult(answer="ok", sources=[]),
    )
    assert client.post("/ask", json={"question": "测试"}).status_code == 200


def test_server_auth_locks_client_after_three_failures(monkeypatch) -> None:
    monkeypatch.setattr(server, "demo_password", lambda: DEMO_TEST_PASSWORD)
    monkeypatch.setattr(
        server,
        "answer_question",
        lambda question, **kwargs: AnswerResult(answer="ok", sources=[]),
    )
    client = TestClient(server.app)
    locked_headers = {"X-Demo-Client-Id": "lock-client-a"}

    res = client.post(
        "/auth/check", json={"password": "wrong"}, headers=locked_headers
    )
    assert res.status_code == 401
    assert "还剩 2 次机会" in res.json()["detail"]
    res = client.post(
        "/auth/check", json={"password": "wrong"}, headers=locked_headers
    )
    assert res.status_code == 401
    assert "还剩 1 次机会" in res.json()["detail"]
    res = client.post(
        "/auth/check", json={"password": "wrong"}, headers=locked_headers
    )
    assert res.status_code == 429

    correct_while_locked = client.post(
        "/auth/check",
        json={"password": DEMO_TEST_PASSWORD},
        headers=locked_headers,
    )
    assert correct_while_locked.status_code == 429

    ask_while_locked = client.post(
        "/ask",
        json={"question": "测试"},
        headers={**AUTH_HEADERS, "X-Demo-Client-Id": "lock-client-a"},
    )
    assert ask_while_locked.status_code == 429

    other_headers = {**AUTH_HEADERS, "X-Demo-Client-Id": "lock-client-b"}
    other = client.post("/ask", json={"question": "测试"}, headers=other_headers)
    assert other.status_code == 200


def test_server_auth_status_reports_lock(monkeypatch) -> None:
    monkeypatch.setattr(server, "demo_password", lambda: DEMO_TEST_PASSWORD)
    client = TestClient(server.app)
    headers = {"X-Demo-Client-Id": "lock-client-status"}
    for _ in range(3):
        client.post("/auth/check", json={"password": "wrong"}, headers=headers)

    data = client.get("/auth/status", headers=headers).json()
    assert data["enabled"] is True
    assert data["locked"] is True
    assert data["retry_after"] > 0


def test_export_pdf_returns_pdf() -> None:
    client = TestClient(server.app)
    response = client.post(
        "/export-pdf",
        headers=AUTH_HEADERS,
        json={
            "question": "退货时间是几天？",
            "answer": "退货时间为 3 天。\n- 要点一\n- 要点二",
            "sources": [RULES_TXT],
        },
    )
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"
    assert response.content.startswith(b"%PDF")
    assert "filename" in response.headers["content-disposition"].lower()
    text = _pdf_text(response.content)
    assert "来源" not in text
    assert RULES_TXT not in text


def test_export_pdf_with_messages_returns_pdf() -> None:
    client = TestClient(server.app)
    response = client.post(
        "/export-pdf",
        headers=AUTH_HEADERS,
        json={
            "messages": [
                {"role": "user", "content": "退货时间是几天？"},
                {
                    "role": "assistant",
                    "content": "退货时间为 3 天。",
                    "sources": [RULES_TXT],
                },
            ]
        },
    )
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"
    assert response.content.startswith(b"%PDF")
    text = _pdf_text(response.content)
    assert "来源" not in text
    assert RULES_TXT not in text


def test_export_pdf_with_messages_skips_unknown() -> None:
    client = TestClient(server.app)
    response = client.post(
        "/export-pdf",
        headers=AUTH_HEADERS,
        json={
            "messages": [
                {"role": "user", "content": "数字孪生是什么？"},
                {"role": "assistant", "content": "数字孪生是物理对象的数字化映射。"},
                {"role": "user", "content": "航母航速是多少？"},
                {
                    "role": "assistant",
                    "content": UNKNOWN_ANSWER,
                    "unknown": True,
                },
            ]
        },
    )
    assert response.status_code == 200
    text = _pdf_text(response.content)
    assert "数字孪生是物理对象的数字化映射" in text
    assert "航母航速是多少" not in text
    assert "我还不懂这些" not in text


def test_export_pdf_without_content_returns_400() -> None:
    client = TestClient(server.app)
    assert client.post("/export-pdf", json={}, headers=AUTH_HEADERS).status_code == 400
    assert (
        client.post(
            "/export-pdf", json={"messages": []}, headers=AUTH_HEADERS
        ).status_code
        == 400
    )


def test_build_qa_pdf_supports_chinese() -> None:
    pdf = build_qa_pdf(
        question="苏州市数据条例对公共数据开放有什么规定？",
        answer="公共数据开放遵循合法、公正、安全原则。\n\n**重点：** 坚持开放共享。",
        sources=["来源文件.pdf"],
    )
    assert pdf.startswith(b"%PDF")
    assert len(pdf) > 1000
    text = _pdf_text(pdf)
    assert "苏州市数据条例对公共数据开放有什么规定" in text
    assert "公共数据开放遵循合法、公正、安全原则" in text
    assert "来源" not in text
    assert "来源文件.pdf" not in text


def test_build_conversation_pdf_supports_multiple_messages() -> None:
    pdf = build_conversation_pdf(
        [
            {"role": "user", "content": "退货时间是几天？"},
            {
                "role": "assistant",
                "content": "退货时间为 3 天。\n- 要点一",
                "sources": [RULES_TXT],
            },
            {"role": "user", "content": "那退款期限呢？"},
            {
                "role": "assistant",
                "content": "退款期限也是 3 天。",
                "sources": [RULES_TXT],
            },
        ]
    )
    assert pdf.startswith(b"%PDF")
    assert len(pdf) > 1000
    text = _pdf_text(pdf)
    assert "退货时间为 3 天" in text
    assert "退款期限也是 3 天" in text
    assert "来源" not in text
    assert RULES_TXT not in text


def test_build_conversation_pdf_skips_unknown_turns() -> None:
    pdf = build_conversation_pdf(
        [
            {"role": "user", "content": "数字孪生是什么？"},
            {"role": "assistant", "content": "数字孪生是物理对象的数字化映射。"},
            {"role": "user", "content": "航母航速是多少？"},
            {"role": "assistant", "content": UNKNOWN_ANSWER, "unknown": True},
            {"role": "user", "content": "园区对数字孪生有什么期望？"},
            {"role": "assistant", "content": "园区期望用于智慧治理。"},
        ]
    )
    text = _pdf_text(pdf)
    assert "数字孪生是物理对象的数字化映射" in text
    assert "园区期望用于智慧治理" in text
    assert "航母航速是多少" not in text
    assert "我还不懂这些" not in text


def test_server_upload(monkeypatch, tmp_path) -> None:
    monkeypatch.setattr(server, "pdf_dir", lambda: tmp_path)
    monkeypatch.setattr(server, "ingest", lambda **kwargs: 3)
    client = TestClient(server.app)

    response = client.post(
        "/upload",
        headers=AUTH_HEADERS,
        files={"file": (ENGLISH_RULES_TXT, "退货时间为 3 天。", "text/plain")},
    )
    assert response.status_code == 200
    body = response.json()
    assert body["added_chunks"] == 3
    assert (tmp_path / body["filename"]).exists()


def test_server_upload_rejects_unsupported_type(monkeypatch, tmp_path) -> None:
    monkeypatch.setattr(server, "pdf_dir", lambda: tmp_path)
    client = TestClient(server.app)

    response = client.post(
        "/upload",
        headers=AUTH_HEADERS,
        files={"file": ("rules.xyz", b"unknown", "application/octet-stream")},
    )
    assert response.status_code == 400
    assert "不支持的文件类型" in response.json()["detail"]


def test_load_docx(tmp_path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("第一章 退货规则", level=1)
    doc.add_paragraph("退货时间为 3 天。")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "项目"
    table.cell(0, 1).text = "期限"
    table.cell(1, 0).text = "退货"
    table.cell(1, 1).text = "3 天"
    path = tmp_path / RULES_DOCX
    doc.save(str(path))

    docs = load_file(path)
    assert len(docs) == 1
    assert "退货时间为 3 天" in docs[0].page_content
    assert "3 天" in docs[0].page_content
    assert docs[0].metadata["source"] == RULES_DOCX


def test_load_docx_with_null_relationship_falls_back(tmp_path) -> None:
    from docx import Document as DocxDocument

    base = tmp_path / "base.docx"
    doc = DocxDocument()
    doc.add_paragraph("退货时间为 3 天。")
    doc.add_paragraph("第二条规则：4 天。")
    doc.save(str(base))

    import zipfile

    with zipfile.ZipFile(base) as archive:
        entries = {name: archive.read(name) for name in archive.namelist()}
    rels = entries["word/_rels/document.xml.rels"].decode("utf-8")
    null_rel = (
        '<Relationship Id="rIdNull" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject" '
        'Target="NULL"/>'
    )
    entries["word/_rels/document.xml.rels"] = rels.replace(
        "</Relationships>", f"{null_rel}</Relationships>"
    ).encode("utf-8")

    path = tmp_path / RULES_DOCX
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, blob in entries.items():
            archive.writestr(name, blob)

    from docx import Document as BrokenDocx

    import pytest

    with pytest.raises(KeyError):
        BrokenDocx(str(path))

    docs = load_file(path)
    assert len(docs) == 1
    assert "退货时间为 3 天" in docs[0].page_content
    assert "第二条规则" in docs[0].page_content


def test_load_doc_disguised_as_docx(tmp_path) -> None:
    from docx import Document

    doc = Document()
    doc.add_paragraph("退货时间为 3 天。")
    path = tmp_path / RULES_DOC
    doc.save(str(path))

    docs = load_file(path)
    assert len(docs) == 1
    assert "退货时间为 3 天" in docs[0].page_content


def test_load_doc_rtf(tmp_path) -> None:
    path = tmp_path / RULES_DOC
    path.write_bytes(b"{\\rtf1\\ansi Return policy: 3 days.}")

    docs = load_file(path)
    assert len(docs) == 1
    assert "Return policy" in docs[0].page_content


def test_load_doc_html(tmp_path) -> None:
    path = tmp_path / RULES_DOC
    path.write_text("<html><body><p>退货时间为 3 天。</p></body></html>", encoding="utf-8")

    docs = load_file(path)
    assert len(docs) == 1
    assert "退货时间为 3 天" in docs[0].page_content


def test_load_doc_ole_via_com(monkeypatch, tmp_path) -> None:
    import shutil

    from docx import Document

    real_doc = Document()
    real_doc.add_paragraph("退货时间为 3 天。")
    docx_path = tmp_path / "converted.docx"
    real_doc.save(str(docx_path))

    def fake_com_save_as(path, app_progids, temp_name, save_format):
        target_dir = tmp_path / "office"
        target_dir.mkdir(exist_ok=True)
        target = target_dir / temp_name
        shutil.copyfile(docx_path, target)
        return target

    monkeypatch.setattr("kb_agent.loaders._com_save_as", fake_com_save_as)
    ole_path = tmp_path / RULES_DOC
    ole_path.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 8)

    docs = load_file(ole_path)
    assert len(docs) == 1
    assert "退货时间为 3 天" in docs[0].page_content
    assert docs[0].metadata["source"] == RULES_DOC
    assert docs[0].metadata["converted_from"] == ".doc"


def test_load_pptx(tmp_path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "退货规则"
    slide.placeholders[1].text = "退货时间为 3 天。"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text = "项目"
    table_shape.table.cell(0, 1).text = "期限"
    table_shape.table.cell(1, 0).text = "退货"
    table_shape.table.cell(1, 1).text = "3 天"
    path = tmp_path / RULES_PPTX
    presentation.save(str(path))

    docs = load_file(path)
    assert len(docs) == 1
    assert "退货时间为 3 天" in docs[0].page_content
    assert "3 天" in docs[0].page_content
    assert docs[0].metadata["slide"] == 1


def test_load_ppt_disguised_as_pptx(tmp_path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "退货规则"
    path = tmp_path / RULES_PPT
    presentation.save(str(path))

    docs = load_file(path)
    assert len(docs) == 1
    assert "退货规则" in docs[0].page_content


def test_load_wps_plain_text(tmp_path) -> None:
    path = tmp_path / RULES_WPS
    path.write_text("退货时间为 3 天。\n", encoding="utf-8")

    docs = load_file(path)
    assert len(docs) == 1
    assert "退货时间为 3 天" in docs[0].page_content


def test_load_doc_without_com_raises_clear_error(monkeypatch, tmp_path) -> None:
    def fail_com(*args, **kwargs):
        raise RuntimeError("没有可用的 Office/WPS")

    monkeypatch.setattr("kb_agent.loaders._com_save_as", fail_com)
    ole_path = tmp_path / RULES_DOC
    ole_path.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 8)

    try:
        load_file(ole_path)
    except ValueError as exc:
        assert "Office/WPS" in str(exc)
    else:
        raise AssertionError("缺少 Office/WPS 时应抛出 ValueError")


def test_load_doc_ole_fallback_text(monkeypatch, tmp_path) -> None:
    def fail_com(*args, **kwargs):
        raise RuntimeError("没有可用的 Office/WPS")

    monkeypatch.setattr("kb_agent.loaders._com_save_as", fail_com)
    monkeypatch.setattr(
        "kb_agent.loaders._extract_ole_doc_text",
        lambda path: "退货时间为 3 天。\n数据条例第二十一条。",
    )
    ole_path = tmp_path / RULES_DOC
    ole_path.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 8)

    docs = load_file(ole_path)
    assert len(docs) == 1
    assert "退货时间为 3 天" in docs[0].page_content
    assert docs[0].metadata["converted_from"] == ".doc"


def test_load_xlsx(tmp_path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "退货表"
    sheet.append(["项目", "期限"])
    sheet.append(["退货", "3 天"])
    path = tmp_path / RULES_XLSX
    workbook.save(str(path))

    docs = load_file(path)
    assert len(docs) == 1
    assert "退货表" in docs[0].page_content
    assert "3 天" in docs[0].page_content
    assert docs[0].metadata["section_path"] == "退货表"


def test_load_csv(tmp_path) -> None:
    path = tmp_path / RULES_CSV
    path.write_text("项目,期限\n退货,3 天\n", encoding="utf-8-sig")

    docs = load_file(path)
    assert len(docs) == 1
    assert "3 天" in docs[0].page_content


def test_text_pdf_loads_without_ocr(tmp_path) -> None:
    import pymupdf

    path = tmp_path / RULES_PDF
    pdf = pymupdf.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Return policy: 3 days.")
    pdf.save(str(path))
    pdf.close()

    docs = load_file(path)
    assert len(docs) == 1
    assert "3 days" in docs[0].page_content
    assert docs[0].metadata.get("ocr") is None


def test_scanned_pdf_falls_back_to_ocr(monkeypatch, tmp_path) -> None:
    import io

    import pymupdf
    from PIL import Image

    path = tmp_path / SCANNED_PDF
    pdf = pymupdf.open()
    page = pdf.new_page(width=300, height=300)
    buffer = io.BytesIO()
    Image.new("RGB", (300, 300), "white").save(buffer, format="PNG")
    page.insert_image(page.rect, stream=buffer.getvalue())
    pdf.save(str(path))
    pdf.close()

    monkeypatch.setattr(
        "kb_agent.loaders._ocr_page",
        lambda page, path, page_number: "扫描件 OCR 文本",
    )
    docs = load_file(path)
    assert len(docs) == 1
    assert "OCR 文本" in docs[0].page_content
    assert docs[0].metadata.get("ocr") is True


def test_load_png_uses_ocr(monkeypatch, tmp_path) -> None:
    from PIL import Image

    path = tmp_path / RULES_PNG
    Image.new("RGB", (300, 300), "white").save(path)

    class _FakeOcrEngine:
        def __call__(self, array):
            result = [([[0, 0], [100, 0], [100, 20], [0, 20]], "退货时间为 3 天", 0.99)]
            return result, 0.01

    monkeypatch.setattr("kb_agent.loaders._get_ocr_engine", lambda: _FakeOcrEngine())
    monkeypatch.setattr("kb_agent.loaders.ocr_cache.cached_text", lambda *a, **k: None)
    monkeypatch.setattr("kb_agent.loaders.ocr_cache.save_text", lambda *a, **k: None)

    docs = load_file(path)
    assert len(docs) == 1
    assert "退货时间为 3 天" in docs[0].page_content
    assert docs[0].metadata.get("ocr") is True
    assert docs[0].metadata["source"] == RULES_PNG


def test_ingest_accepts_multiple_formats(tmp_path) -> None:
    from docx import Document
    from openpyxl import Workbook

    source_dir = tmp_path / "pdfs"
    source_dir.mkdir()
    (source_dir / RULES_TXT).write_text("退货时间为 3 天。\n", encoding="utf-8")

    docx = Document()
    docx.add_paragraph("退货时间规则：3 天。")
    docx.save(str(source_dir / RULES_DOCX))

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "退货表"
    sheet.append(["项目", "期限"])
    sheet.append(["退货", "3 天"])
    workbook.save(str(source_dir / RULES_XLSX))

    index_dir = tmp_path / "index"
    added = ingest(
        directory=source_dir,
        index_directory=index_dir,
        embedding=HashEmbeddings(),
        rebuild=True,
    )
    assert added == 3

    vectorstore = load_vectorstore(index_directory=index_dir, embedding=HashEmbeddings())
    retriever = build_retriever(vectorstore, top_k=3, score_threshold=0.0)
    documents = retriever.invoke("退货时间是几天？")
    assert any("3 天" in doc.page_content for doc in documents)


def test_structure_hits_prioritize_specific_section() -> None:
    from types import SimpleNamespace

    plan_docs = [
        Document(
            page_content=f"章节 {index}",
            metadata={
                "source": PLAN_PDF,
                "section_path": f"三、重点行动 > （{index}）数据要素×场景",
            },
        )
        for index in range(20)
    ]
    docx_doc = Document(
        page_content="第二章 认定条件",
        metadata={"source": MEASURE_DOCX, "section_path": "第二章 认定条件"},
    )
    all_docs = plan_docs + [docx_doc]
    fake_vectorstore = SimpleNamespace(
        docstore=SimpleNamespace(_dict={id(doc): doc for doc in all_docs})
    )

    class FakeEnsemble:
        def invoke(self, question: str) -> list[Document]:
            return []

    retriever = GroundedRetriever(FakeEnsemble(), fake_vectorstore, 0.0)
    hits = retriever._structure_hits(
        "江苏省数据企业评估认定管理办法中平台企业的认定条件是什么"
    )
    assert hits
    assert hits[0].metadata["source"] == MEASURE_DOCX


def test_tier_metadata_through_ingest(tmp_path) -> None:
    source_dir = tmp_path / "pdfs"
    tier_dir = source_dir / "01国家政策"
    tier_dir.mkdir(parents=True)
    (tier_dir / ENGLISH_RULES_TXT).write_text("退货时间为 3 天。\n", encoding="utf-8")
    index_dir = tmp_path / "index"

    added = ingest(
        directory=source_dir,
        index_directory=index_dir,
        embedding=HashEmbeddings(),
        rebuild=True,
    )
    assert added == 1
    vectorstore = load_vectorstore(index_directory=index_dir, embedding=HashEmbeddings())
    docs = list(vectorstore.docstore._dict.values())
    assert docs[0].metadata.get("tier") == "国家政策"


def test_detect_tier() -> None:
    assert detect_tier("国家政策里的数据要素×行动计划是什么") == "国家政策"
    assert detect_tier("全省数据工作会议报告主要内容") == "省级政策"
    assert detect_tier("苏州的市级政策有哪些") == "市级政策"
    assert detect_tier("园区政策中有什么要求") == "园区政策"
    assert detect_tier("其他地区政策如何规定") == "其他地区政策"
    assert detect_tier("退货时间是几天") is None
    # 外省/外地文件可能暂存在非对应层级目录，问具体地区时跨层级检索
    assert detect_tier("山东省省级语料券奖补资金管理办法") is None
    assert detect_tier("济南市数据要素奖补资金实施细则") is None
    assert detect_tier("深圳市人工智能语料券专项资金操作规程") is None
    assert detect_tier("上海市进一步扩大人工智能应用的若干措施") is None


def test_other_region_source_markers() -> None:
    assert _is_other_region_source("上海市数据条例.txt") is True
    assert _is_other_region_source("青岛市数据要素市场发展政策实施细则.pdf") is True
    assert _is_other_region_source("深政数规〔2025〕1号.pdf") is True
    assert _is_other_region_source("苏州市数据条例.doc") is False
    assert _is_other_region_source("江苏省数据条例.doc") is False
    assert _is_other_region_source("无锡市数据条例.docx") is False
    assert _is_other_region_source("玄武区关于促进数据产业发展.doc") is False


def test_local_source_names_keeps_jiangsu_docs_in_other_tier() -> None:
    docs = [
        Document(
            page_content="杭州高新区关于促进数据产业高质量发展的实施意见。",
            metadata={"source": "关于促进数据产业高质量发展的实施意见.txt", "tier": "其他地区政策"},
        ),
        Document(
            page_content="天聚地合（苏州）科技股份有限公司可信数据空间项目。",
            metadata={"source": "可信数据空间.txt", "tier": "其他地区政策"},
        ),
        Document(
            page_content="上海市数据条例内容。",
            metadata={"source": "上海市数据条例.txt", "tier": "其他地区政策"},
        ),
        Document(
            page_content="国家数据要素政策内容。",
            metadata={"source": "国家数据政策.txt", "tier": "国家政策"},
        ),
    ]
    local = _local_source_names(docs)
    assert "关于促进数据产业高质量发展的实施意见.txt" not in local
    assert "上海市数据条例.txt" not in local
    assert "可信数据空间.txt" in local
    assert "国家数据政策.txt" in local


def test_needs_other_region_docs() -> None:
    assert _needs_other_region_docs("苏州市数据条例有哪些规定") is False
    assert _needs_other_region_docs("数据要素政策有哪些") is False
    assert _needs_other_region_docs("上海市数据条例有哪些规定") is True
    assert _needs_other_region_docs("苏州和上海有什么区别") is True
    assert _needs_other_region_docs("对比一下数据产业政策") is False
    assert _needs_other_region_docs("语料券奖补资金怎么申请") is True
    assert _needs_other_region_docs("其他地区政策有哪些") is True


def test_default_retrieval_excludes_other_region_docs(tmp_path) -> None:
    source_dir = tmp_path / "pdfs"
    source_dir.mkdir()
    national_dir = source_dir / "01国家政策"
    other_dir = source_dir / "05其他地区政策"
    national_dir.mkdir()
    other_dir.mkdir()
    (national_dir / "国家数据政策.txt").write_text(
        "国家数据要素政策内容。", encoding="utf-8"
    )
    (other_dir / "上海市数据条例.txt").write_text(
        "上海市数据条例内容。", encoding="utf-8"
    )
    (other_dir / "关于促进数据产业高质量发展的实施意见.txt").write_text(
        "杭州高新区关于促进数据产业高质量发展的实施意见：数据要素政策内容。",
        encoding="utf-8",
    )
    index_dir = tmp_path / "index"
    ingest(
        directory=source_dir,
        index_directory=index_dir,
        embedding=HashEmbeddings(),
        rebuild=True,
    )
    vectorstore = load_vectorstore(index_directory=index_dir, embedding=HashEmbeddings())
    result = answer_question(
        "数据要素政策有哪些",
        vectorstore=vectorstore,
        llm=StubChatModel(answer="国家数据要素政策。"),
        score_threshold=0.0,
    )
    assert result.sources == ["国家数据政策.txt"]
    assert all("上海市" not in source for source in result.sources)
    assert all("促进数据产业高质量发展" not in source for source in result.sources)


def test_explicit_other_region_question_includes_other_region_docs(tmp_path) -> None:
    source_dir = tmp_path / "pdfs"
    source_dir.mkdir()
    national_dir = source_dir / "01国家政策"
    other_dir = source_dir / "05其他地区政策"
    national_dir.mkdir()
    other_dir.mkdir()
    (national_dir / "国家数据政策.txt").write_text(
        "国家数据要素政策内容。", encoding="utf-8"
    )
    (other_dir / "上海市数据条例.txt").write_text(
        "上海市数据条例内容。", encoding="utf-8"
    )
    index_dir = tmp_path / "index"
    ingest(
        directory=source_dir,
        index_directory=index_dir,
        embedding=HashEmbeddings(),
        rebuild=True,
    )
    vectorstore = load_vectorstore(index_directory=index_dir, embedding=HashEmbeddings())
    result = answer_question(
        "上海市数据条例有哪些规定",
        vectorstore=vectorstore,
        llm=StubChatModel(answer="上海市数据条例内容。"),
        score_threshold=0.0,
    )
    assert "上海市数据条例.txt" in result.sources


def test_retriever_tier_filter() -> None:
    from types import SimpleNamespace

    national = Document(
        page_content="国家数据要素行动计划",
        metadata={"source": NATIONAL_TXT, "tier": "国家政策", "section_path": ""},
    )
    provincial = Document(
        page_content="省级数据企业管理办法",
        metadata={"source": PROVINCIAL_DOCX, "tier": "省级政策", "section_path": ""},
    )
    municipal = Document(
        page_content="市级可信数据空间方案",
        metadata={"source": MUNICIPAL_DOCX, "tier": "市级政策", "section_path": ""},
    )
    all_docs = [national, provincial, municipal]
    fake_vectorstore = SimpleNamespace(
        docstore=SimpleNamespace(_dict={id(doc): doc for doc in all_docs})
    )
    fake_vectorstore.similarity_search_with_relevance_scores = (
        lambda question, k: [(doc, 0.6) for doc in all_docs]
    )

    class StaticKeyword:
        def invoke(self, question: str) -> list[Document]:
            return [provincial]

    retriever = GroundedRetriever(None, fake_vectorstore, 0.0, tier="省级政策", top_k=5)
    retriever.keyword = StaticKeyword()
    docs = retriever.invoke("数据企业管理办法")
    assert docs
    assert all(doc.metadata["tier"] == "省级政策" for doc in docs)


def test_empty_tier_returns_unknown(tmp_path) -> None:
    source_dir = tmp_path / "pdfs"
    source_dir.mkdir()
    tier_dir = source_dir / "01国家政策"
    tier_dir.mkdir()
    (tier_dir / ENGLISH_RULES_TXT).write_text("退货时间为 3 天。", encoding="utf-8")
    index_dir = tmp_path / "index"
    ingest(
        directory=source_dir,
        index_directory=index_dir,
        embedding=HashEmbeddings(),
        rebuild=True,
    )
    vectorstore = load_vectorstore(index_directory=index_dir, embedding=HashEmbeddings())
    retriever = build_retriever(
        vectorstore,
        top_k=5,
        score_threshold=0.0,
        tier="园区政策",
    )

    assert retriever.invoke("园区政策有什么") == []
    result = answer_question(
        "园区政策有什么",
        vectorstore=vectorstore,
        retriever=retriever,
        llm=StubChatModel(),
    )
    assert result.answer == UNKNOWN_ANSWER_WITH_OFFER
    assert result.sources == []
    assert result.unknown is True


def test_chunk_includes_source_label() -> None:
    doc = Document(
        page_content="退货时间为 3 天。",
        metadata={"source": RULES_TXT},
    )
    chunks = split_documents([doc])
    assert chunks
    assert f"【来源：{RULES_TXT}】" in chunks[0].page_content


def test_detect_target_sources() -> None:
    docs = [
        Document(page_content="x", metadata={"source": FEEDBACK_DOC}),
        Document(page_content="y", metadata={"source": FEEDBACK_0116_DOC}),
        Document(page_content="z", metadata={"source": PLAN_PDF}),
    ]
    assert sorted(detect_target_sources(f"检查一下《{FEEDBACK_DOC.split('.')[0]}》有没有答案", docs)) == [
        FEEDBACK_DOC,
        FEEDBACK_0116_DOC,
    ]
    assert detect_target_sources("修改建议有哪些", docs) == []


def test_target_source_scope_retrieval() -> None:
    from types import SimpleNamespace

    target = Document(
        page_content="大数据局反馈 修改建议：增加高质量数据集表述。",
        metadata={"source": FEEDBACK_DOC, "section_path": ""},
    )
    other = Document(
        page_content="行动计划内容。",
        metadata={"source": PLAN_PDF, "section_path": "三、重点行动"},
    )
    all_docs = [target, other]
    fake_vectorstore = SimpleNamespace(
        docstore=SimpleNamespace(_dict={id(doc): doc for doc in all_docs})
    )
    fake_vectorstore.similarity_search_with_relevance_scores = (
        lambda question, k: [(doc, 0.6) for doc in all_docs]
    )

    retriever = build_retriever(
        fake_vectorstore,
        top_k=5,
        score_threshold=0.0,
        target_sources=[FEEDBACK_DOC],
    )
    docs = retriever.invoke("修改建议有哪些")
    assert docs
    assert all(doc.metadata["source"] == FEEDBACK_DOC for doc in docs)


def test_structure_hits_do_not_drown_normal_hits() -> None:
    from types import SimpleNamespace

    target = Document(
        page_content="大数据局反馈 修改建议：增加高质量数据集表述。",
        metadata={"source": FEEDBACK_DOC, "section_path": ""},
    )
    noisy = [
        Document(
            page_content=f"章节 {index}",
            metadata={
                "source": PLAN_PDF,
                "section_path": "三、数据要素行动计划 > （四）数据要素×场景",
            },
        )
        for index in range(30)
    ]
    all_docs = [target] + noisy
    fake_vectorstore = SimpleNamespace(
        docstore=SimpleNamespace(_dict={id(doc): doc for doc in all_docs})
    )
    fake_vectorstore.similarity_search_with_relevance_scores = (
        lambda question, k: [(doc, 0.6) for doc in all_docs]
    )

    class FakeEnsemble:
        def invoke(self, question):
            return [target] + noisy[:20]

    retriever = GroundedRetriever(FakeEnsemble(), fake_vectorstore, 0.0)
    docs = retriever.invoke("大数据局给了什么修改建议")
    assert any(doc.metadata["source"] == FEEDBACK_DOC for doc in docs[:12])


def test_is_summary_question_detects_summary_keywords() -> None:
    assert _is_summary_question("总结一下国家政策有哪些重点行动")
    assert _is_summary_question("梳理市级政策中数据产业的支持措施")
    assert not _is_summary_question("退货时间是几天？")


def test_is_report_question_detects_report_keywords() -> None:
    assert _is_report_question("帮我出一份《苏州市数据条例》的政策解读报告")
    assert _is_report_question("解读一下苏州工业园区数字贸易三年行动计划")
    assert _is_report_question("汇报一下江苏省数据条例的要点")
    assert _is_report_question("这份文件的主要内容是什么")
    assert not _is_report_question("退货时间是几天？")
    assert not _is_report_question("苏州市数据条例对公共数据开放有什么规定")


def test_report_prompt_grounding_rules() -> None:
    assert "结构化解读/报告" in REPORT_SYSTEM_PROMPT
    assert "只能填写知识库内容中明确写出的信息" in REPORT_SYSTEM_PROMPT
    assert "禁止编造" in REPORT_SYSTEM_PROMPT
    assert "不得添加知识库之外的任何信息" in REPORT_SYSTEM_PROMPT
    assert "（《文件名》）" in REPORT_SYSTEM_PROMPT
    assert "知识库未提及/未列出/缺失" in REPORT_SYSTEM_PROMPT
    assert "为了凑齐模板而补内容" in REPORT_SYSTEM_PROMPT
    assert "不要要求用户上传或补充文档" in REPORT_SYSTEM_PROMPT


def test_answer_question_uses_report_prompt() -> None:
    documents = [
        Document(
            page_content="苏州市数据条例：公共数据开放遵循合法、公正、安全原则。",
            metadata={"source": "苏州市数据条例.docx", "page": 1},
        )
    ]
    model = RecordingChatModel(answer="结构化报告正文")
    result = answer_question(
        "请解读《苏州市数据条例》并生成一份结构化报告",
        retriever=StaticRetriever(documents),
        llm=model,
    )
    assert result.answer == "结构化报告正文"
    assert result.sources == ["苏州市数据条例.docx"]
    assert "结构化解读/报告" in model.last_messages[0].content
    assert "只能填写知识库内容中明确写出的信息" in model.last_messages[0].content


def test_report_mode_hides_sources_when_unknown() -> None:
    documents = [
        Document(
            page_content="苏州市数据条例：公共数据开放遵循合法、公正、安全原则。",
            metadata={"source": "苏州市数据条例.docx", "page": 1},
        )
    ]
    result = answer_question(
        "请解读《苏州市数据条例》",
        retriever=StaticRetriever(documents),
        llm=StubChatModel(answer=UNKNOWN_ANSWER),
    )
    assert result.answer == UNKNOWN_ANSWER_WITH_OFFER
    assert result.sources == []
    assert result.unknown is True


def test_complete_source_context_groups_and_completes_sources() -> None:
    docs = [
        Document(page_content=f"A 第{page}条", metadata={"source": "A.pdf", "page": page})
        for page in (3, 1, 2)
    ] + [
        Document(page_content=f"B 第{page}条", metadata={"source": "B.pdf", "page": page})
        for page in (2, 1)
    ]
    hits = [docs[0], docs[3]]
    result = _complete_source_context(hits, docs, limit=10, full_source_max=80)
    assert len(result) == 5
    assert [doc.metadata["source"] for doc in result] == ["A.pdf"] * 3 + ["B.pdf"] * 2
    assert [doc.metadata["page"] for doc in result[:3]] == [1, 2, 3]


def test_complete_source_context_prioritizes_source_completeness() -> None:
    docs = [
        Document(page_content=f"A {page}", metadata={"source": "A.pdf", "page": page})
        for page in (1, 2, 3)
    ] + [
        Document(page_content=f"B {page}", metadata={"source": "B.pdf", "page": page})
        for page in (1, 2)
    ]
    noise = [
        Document(page_content=f"noise {index}", metadata={"source": f"noise{index}.pdf"})
        for index in range(50)
    ]
    hits = [docs[1], docs[3]] + noise
    result = _complete_source_context(hits, docs + noise, limit=10, full_source_max=80)
    sources = [doc.metadata["source"] for doc in result]
    assert sources == ["A.pdf"] * 3 + ["B.pdf"] * 2 + [f"noise{i}.pdf" for i in range(5)]


def test_system_prompt_requires_per_source_answer() -> None:
    assert "涉及多份文件时，先归纳共同的核心要点" in SYSTEM_PROMPT
    assert "确需区分文件时再按“《文件名》”分组" in SYSTEM_PROMPT
    assert "禁止把不同文件的条目混成一份清单" in SYSTEM_PROMPT
    assert "不要在回答中评价知识库是否完整" in SYSTEM_PROMPT


def test_system_prompt_uses_euphemistic_unknown_and_concise_summary() -> None:
    assert "我还不懂这些" in UNKNOWN_ANSWER
    assert "换个问法" in UNKNOWN_ANSWER
    assert "补充相关文档" not in UNKNOWN_ANSWER
    assert "上传" not in UNKNOWN_ANSWER
    assert "长度要适中" in SYSTEM_PROMPT
    assert "概括核心要点" in SYSTEM_PROMPT
    assert "不得添加知识库之外的任何信息" in SYSTEM_PROMPT
    assert "不要逐文件复制目标、任务和全部细项" in SYSTEM_PROMPT
    assert "（《文件名》）" in SYSTEM_PROMPT
    assert "不要要求用户上传或补充文档" in SYSTEM_PROMPT


def test_sources_list_documents_without_page_numbers() -> None:
    docs = [
        Document(page_content="a", metadata={"source": "A.pdf", "page": 1}),
        Document(page_content="b", metadata={"source": "A.pdf", "page": 2}),
        Document(page_content="c", metadata={"source": "B.pdf", "page": 1}),
    ]
    assert _format_sources(docs) == ["A.pdf", "B.pdf"]


def test_load_documents_skips_office_temp_files(tmp_path) -> None:
    source_dir = tmp_path / "pdfs"
    source_dir.mkdir()
    (source_dir / REAL_TXT).write_text("真实内容", encoding="utf-8")
    (source_dir / TEMP_DOCX).write_text("junk", encoding="utf-8")
    docs = load_documents(source_dir)
    assert len(docs) == 1
    assert docs[0].metadata["source"] == REAL_TXT
