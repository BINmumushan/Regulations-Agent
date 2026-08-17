"""多格式文档读取：PDF、Office/文本/表格、WPS、演示文稿。"""

from __future__ import annotations

import csv
import html
from html.parser import HTMLParser
import io
import logging
import re
import shutil
import tempfile
import threading
from datetime import date, datetime
from pathlib import Path

from langchain_core.documents import Document

from . import ocr_cache
from .config import get_env, get_int_env, pdf_dir

logger = logging.getLogger(__name__)

SUPPORTED_SUFFIXES = {
    ".pdf",
    ".jpg",
    ".jpeg",
    ".png",
    ".txt",
    ".md",
    ".csv",
    ".docx",
    ".doc",
    ".ppt",
    ".pptx",
    ".wps",
    ".xlsx",
    ".xlsm",
    ".xls",
}

_MIN_PDF_TEXT_LEN = 8

_IMAGE_OCR_STRIP = 800
_IMAGE_OCR_OVERLAP = 80
_IMAGE_OCR_DPI = 0  # 图片没有 DPI 概念，固定值仅用于把图片缓存与 PDF 页面缓存隔开

_ocr_engine = None
_ocr_lock = threading.Lock()

_TIER_DIR_PREFIX = re.compile(r"^\s*\d+\s*[-_　]?\s*")


def _tier_from_path(path: Path) -> str:
    """从目录名识别政策层级，如 01国家政策 -> 国家政策。"""
    parent = path.parent.name.strip()
    if path.parent.resolve() == pdf_dir().resolve():
        return ""
    return _TIER_DIR_PREFIX.sub("", parent).strip()


def _read_text_bytes(raw: bytes) -> str:
    """按常见中文编码顺序解码文本文件，避免 GBK 等编码的文件乱码。"""
    for encoding in ("utf-8-sig", "gb18030", "utf-16"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("latin-1", errors="replace")


def _text_document(path: Path, content: str, **metadata: object) -> Document:
    meta = {"source": path.name}
    meta.update(metadata)
    return Document(page_content=content.strip(), metadata=meta)


def _load_text_file(path: Path) -> list[Document]:
    return [_text_document(path, _read_text_bytes(path.read_bytes()))]


def _load_csv(path: Path) -> list[Document]:
    text = _read_text_bytes(path.read_bytes())
    lines: list[str] = []
    for row in csv.reader(io.StringIO(text)):
        cells = [cell.strip() for cell in row if cell and cell.strip()]
        if cells:
            lines.append(" | ".join(cells))
    return [_text_document(path, "\n".join(lines))]


def _ocr_enabled() -> bool:
    return get_env("PDF_OCR_ENABLED", "1").strip().lower() not in {"0", "false", "no", "off"}


def _ocr_available() -> bool:
    try:
        import rapidocr_onnxruntime  # noqa: F401

        return True
    except ImportError:
        return False


def _get_ocr_engine():
    global _ocr_engine
    if _ocr_engine is None:
        with _ocr_lock:
            if _ocr_engine is None:
                from rapidocr_onnxruntime import RapidOCR

                _ocr_engine = RapidOCR()
    return _ocr_engine


def _ocr_page(page, path: Path, page_number: int) -> str:
    """对没有文字层的 PDF 页做 OCR，返回按阅读顺序拼接的文本。"""
    dpi = get_int_env("PDF_OCR_DPI", 200)
    cached = ocr_cache.cached_text(path, page_number, dpi)
    if cached is not None:
        return cached
    if not _ocr_enabled():
        return ""
    if not _ocr_available():
        logger.warning(
            "检测到扫描 PDF 页，但未安装 rapidocr-onnxruntime，请先执行 pip install -r requirements.txt"
        )
        return ""
    try:
        import numpy as np
        from PIL import Image

        engine = _get_ocr_engine()
        pix = page.get_pixmap(dpi=get_int_env("PDF_OCR_DPI", 200), alpha=False)
        image = pix.pil_image().convert("RGB")
        result, _ = engine(np.asarray(image)[:, :, ::-1])
    except Exception:
        logger.exception("OCR 识别失败，该页将不进入知识库。")
        return ""
    if not result:
        return ""
    text = ocr_cache.format_ocr_result(result)
    ocr_cache.save_text(path, page_number, text, dpi)
    return text


def _load_pdf(path: Path) -> list[Document]:
    import pymupdf

    docs: list[Document] = []
    with pymupdf.open(str(path)) as pdf:
        scanned_pages = [
            page_number
            for page_number, page in enumerate(pdf, start=1)
            if len(page.get_text("text").strip()) < _MIN_PDF_TEXT_LEN
        ]
        if len(scanned_pages) >= 10 and _ocr_enabled() and _ocr_available():
            try:
                ocr_cache.ocr_pdf_parallel(path)
            except Exception:
                logger.exception("并行 OCR 失败，将逐页识别: %s", path.name)
        for page_number, page in enumerate(pdf, start=1):
            text = page.get_text("text").strip()
            ocr_used = False
            if len(text) < _MIN_PDF_TEXT_LEN:
                ocr_text = _ocr_page(page, path, page_number)
                if ocr_text.strip():
                    text = ocr_text.strip()
                    ocr_used = True
            metadata: dict[str, object] = {"source": path.name, "page": page_number}
            if ocr_used:
                metadata["ocr"] = True
            docs.append(Document(page_content=text, metadata=metadata))
    return docs


def _ocr_image_lines(image) -> list[str]:
    """对图片做 OCR；长图按垂直条带识别并合并，返回按阅读顺序的文本行。"""
    import numpy as np

    engine = _get_ocr_engine()
    width, height = image.size
    lines: list[str] = []
    step = _IMAGE_OCR_STRIP - _IMAGE_OCR_OVERLAP
    for top in range(0, height, step):
        bottom = min(top + _IMAGE_OCR_STRIP, height)
        crop = image.crop((0, top, width, bottom))
        result, _ = engine(np.asarray(crop)[:, :, ::-1])
        for _, text, _score in result or []:
            text = str(text).strip()
            # 相邻条带重叠区会重复识别同一行，连续重复只保留一次
            if not text:
                continue
            if not lines or text != lines[-1]:
                lines.append(text)
    return lines


def _load_image(path: Path) -> list[Document]:
    """对 JPG/PNG 等图片做 OCR，识别文本入库（结果按图片缓存复用）。"""
    from PIL import Image

    cached = ocr_cache.cached_text(path, 1, _IMAGE_OCR_DPI)
    if cached is not None:
        return [_text_document(path, cached, ocr=True)]
    if not _ocr_enabled():
        logger.warning("OCR 已禁用，图片无法识别: %s", path.name)
        return [_text_document(path, "")]
    if not _ocr_available():
        logger.warning(
            "检测到图片，但未安装 rapidocr-onnxruntime，请先执行 pip install -r requirements.txt"
        )
        return [_text_document(path, "")]
    try:
        with Image.open(path) as opened:
            text = "\n".join(_ocr_image_lines(opened.convert("RGB")))
    except Exception:
        logger.exception("图片 OCR 失败，该图片将不进入知识库: %s", path.name)
        return [_text_document(path, "")]
    ocr_cache.save_text(path, 1, text, _IMAGE_OCR_DPI)
    return [_text_document(path, text, ocr=True)]


def _table_row_text(row) -> str:
    cells: list[str] = []
    previous = ""
    for cell in row.cells:
        text = cell.text.strip()
        if text == previous:
            continue
        cells.append(text)
        previous = text
    return " | ".join(cells)


def _load_docx(path: Path) -> list[Document]:
    from docx import Document as DocxDocument
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    try:
        doc = DocxDocument(str(path))
    except KeyError as exc:
        # WPS 等工具生成的 docx 可能带 Target="NULL" 的伪关系，
        # python-docx 会尝试读取不存在的 part 并抛 KeyError，直接走 XML 文本提取。
        logger.warning("python-docx 无法解析（%s），回退到 XML 文本提取: %s", exc, path.name)
        return _load_docx_xml(path)
    lines: list[str] = []
    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            text = Paragraph(child, doc).text.strip()
        elif child.tag == qn("w:tbl"):
            table = Table(child, doc)
            text = "\n".join(_table_row_text(row) for row in table.rows)
        else:
            continue
        if text:
            lines.append(text)
    return [_text_document(path, "\n".join(lines))]


def _load_docx_xml(path: Path) -> list[Document]:
    """不依赖 python-docx，直接从 word/document.xml 提取正文与表格。"""
    import zipfile
    from xml.etree import ElementTree as ET

    w = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"

    def xml_text(element) -> str:
        return "".join(text.text or "" for text in element.iter(w + "t")).strip()

    with zipfile.ZipFile(path) as archive:
        root = ET.fromstring(archive.read("word/document.xml"))
    body = root.find(w + "body")
    lines: list[str] = []
    for child in list(body) if body is not None else []:
        if child.tag == w + "p":
            text = xml_text(child)
            if text:
                lines.append(text)
        elif child.tag == w + "tbl":
            rows: list[str] = []
            for row in child.iter(w + "tr"):
                cells = [xml_text(cell) for cell in row.iter(w + "tc")]
                cells = [cell for cell in cells if cell]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                lines.append("\n".join(rows))
    return [_text_document(path, "\n".join(lines))]


def _sniff_format(path: Path) -> str:
    """根据文件头判断真实格式，避免被扩展名误导。"""
    with path.open("rb") as handle:
        head = handle.read(16)
    if head.startswith(b"\xd0\xcf\x11\xe0"):
        return "ole"
    if head.startswith((b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08")):
        return "zip"
    lowered = head.lstrip().lower()
    if lowered.startswith(b"{\\rtf"):
        return "rtf"
    if lowered.startswith((b"<html", b"<!doctype html", b"<head", b"<body")):
        return "html"
    return "text"


def _load_zip_document(path: Path) -> list[Document]:
    """扩展名为 .doc/.ppt/.wps 但实际是 OOXML 压缩包时，按内部结构解析。"""
    import zipfile

    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
    if "word/document.xml" in names:
        return _load_docx(path)
    if "ppt/presentation.xml" in names:
        return _load_pptx(path)
    raise ValueError(f"无法识别的 OOXML 文档结构: {path.name}")


class _HTMLTextParser(HTMLParser):
    """提取 HTML 正文文本，块级标签处保留换行。"""

    BLOCK_TAGS = {
        "body",
        "br",
        "div",
        "h1",
        "h2",
        "h3",
        "h4",
        "h5",
        "h6",
        "li",
        "p",
        "section",
        "table",
        "td",
        "th",
        "tr",
        "ul",
        "ol",
    }

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in self.BLOCK_TAGS:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self.BLOCK_TAGS:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if data:
            self.parts.append(data)


def _strip_html(content: str) -> str:
    parser = _HTMLTextParser()
    parser.feed(content)
    lines = [line.strip() for line in "".join(parser.parts).splitlines() if line.strip()]
    return "\n".join(lines)


def _load_rtf(path: Path) -> list[Document]:
    from striprtf.striprtf import rtf_to_text

    raw_text = path.read_bytes().decode("gb18030", errors="replace")
    try:
        text = rtf_to_text(raw_text)
    except Exception:
        logger.warning("RTF 解析失败，回退到去除控制字符: %s", path.name)
        text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", raw_text)
    return [_text_document(path, text)]


def _load_plain_office_text(path: Path, fmt: str) -> list[Document]:
    if fmt == "rtf":
        return _load_rtf(path)
    if fmt == "html":
        return [_text_document(path, _strip_html(_read_text_bytes(path.read_bytes())))]
    return _load_text_file(path)


def _com_release_document(document) -> None:
    if document is None:
        return
    try:
        document.Close(False)
    except Exception:
        pass


def _com_quit_app(app) -> None:
    if app is None:
        return
    try:
        app.Quit()
    except Exception:
        pass


def _com_save_as(
    path: Path,
    app_progids: tuple[str, ...],
    temp_name: str,
    save_format: int,
) -> Path:
    """通过 Office/WPS COM 打开旧格式文档，另存为 OOXML 临时文件。"""
    import pythoncom
    import win32com.client
    import time

    com_initialized = False
    try:
        pythoncom.CoInitialize()
        com_initialized = True
    except Exception:
        pass
    temp_dir = Path(tempfile.mkdtemp(prefix="kb_office_"))
    app = None
    document = None
    errors: list[str] = []
    result: Path | None = None
    try:
        for progid in app_progids:
            for attempt in range(2):
                try:
                    temp_path = temp_dir / (temp_name if attempt == 0 else f"retry{attempt}_{temp_name}")
                    try:
                        app = win32com.client.DispatchEx(progid)
                    except Exception:
                        app = win32com.client.Dispatch(progid)
                    try:
                        app.Visible = False
                    except Exception:
                        pass
                    try:
                        app.DisplayAlerts = 0
                    except Exception:
                        pass
                    if progid.upper().startswith(("POWERPOINT", "KWPP", "WPP")):
                        document = app.Presentations.Open(
                            str(path.resolve()), ReadOnly=True, Untitled=False, WithWindow=False
                        )
                        document.SaveAs(str(temp_path), save_format)
                    else:
                        document = app.Documents.Open(
                            str(path.resolve()), ReadOnly=True, AddToRecentFiles=False
                        )
                        try:
                            document.SaveAs2(str(temp_path), FileFormat=save_format)
                        except Exception:
                            document.SaveAs(str(temp_path), FileFormat=save_format)
                    result = temp_path
                    break
                except Exception as exc:
                    errors.append(f"{progid} 第{attempt + 1}次: {exc}")
                    _com_release_document(document)
                    document = None
                    _com_quit_app(app)
                    app = None
                    if attempt == 0:
                        time.sleep(1.0)
            if result is not None:
                break
        if result is None:
            raise RuntimeError(
                f"无法用 Office/WPS COM 打开 {path.name}，请确认已安装 Microsoft Office 或 WPS。"
                f"详细错误: {' | '.join(errors)}"
            )
        return result
    finally:
        _com_release_document(document)
        _com_quit_app(app)
        if com_initialized:
            pythoncom.CoUninitialize()
        if result is None:
            shutil.rmtree(temp_dir, ignore_errors=True)


def _decode_doc_text(raw: bytes) -> str:
    """按 Word 97+ 常见编码尝试解码正文，优先 UTF-16LE，其次 GB18030。"""
    for encoding in ("utf-16-le", "gb18030"):
        try:
            text = raw.decode(encoding)
        except UnicodeDecodeError:
            continue
        printable = sum(1 for ch in text if ch.isprintable() or ch in "\r\n\t")
        if text and printable / len(text) >= 0.5:
            return text
    return raw.decode("gb18030", errors="ignore")


def _clean_ole_text(text: str) -> str:
    """清理 Word 二进制正文中的段落/单元格标记，只保留可读文本。"""
    text = text.replace("\x07", " | ").replace("\x08", "\n")
    cleaned = [
        "\n"
        if char in "\r\n"
        else ("" if ord(char) < 32 and char not in "\t" else char)
        for char in text
    ]
    joined = "".join(cleaned)
    while "\r\n" in joined or "\r\r" in joined or "\n\n\n" in joined:
        joined = joined.replace("\r\n", "\n").replace("\r\r", "\n").replace("\n\n\n", "\n\n")
    joined = re.sub(r"[─\-–—\s]*PAGE\s*\d+[─\-–—\s]*", "", joined, flags=re.IGNORECASE)
    joined = re.sub(r"[ \t]{2,}", " ", joined)
    return joined.strip()


def _extract_ole_doc_text(path: Path) -> str:
    """不依赖 Office/WPS，解析 OLE 二进制 Word 文档的段表并提取正文。"""
    import olefile

    with olefile.OleFileIO(str(path)) as ole:
        if not ole.exists("WordDocument"):
            raise ValueError("OLE 文件中没有 WordDocument 流，可能不是 Word/WPS 文档。")
        word_stream = ole.openstream("WordDocument").read()
        flags = int.from_bytes(word_stream[0x0A:0x0C], "little")
        table_name = "1Table" if flags & 0x0200 else "0Table"

        if not ole.exists(table_name):
            fc_min = int.from_bytes(word_stream[0x18:0x1C], "little")
            lcb_min = int.from_bytes(word_stream[0x1C:0x20], "little")
            return _clean_ole_text(_decode_doc_text(word_stream[fc_min : fc_min + lcb_min]))

        fc_clx = int.from_bytes(word_stream[0x01A2:0x01A6], "little")
        lcb_clx = int.from_bytes(word_stream[0x01A6:0x01AA], "little")
        table_stream = ole.openstream(table_name).read()
        clx = table_stream[fc_clx : fc_clx + lcb_clx]

        pos = 0
        while pos < len(clx) and clx[pos] == 0x01:
            cb = int.from_bytes(clx[pos + 1 : pos + 3], "little")
            pos += 3 + cb
        if pos >= len(clx) or clx[pos] != 0x02:
            raise ValueError("无法在 OLE 表中定位 Word 段表（CLX）。")
        lcb = int.from_bytes(clx[pos + 1 : pos + 5], "little")
        plc = clx[pos + 5 : pos + 5 + lcb]
        if len(plc) < 4 or (len(plc) - 4) % 12 != 0:
            raise ValueError("Word 段表长度异常，无法解析。")
        piece_count = (len(plc) - 4) // 12
        cps = [
            int.from_bytes(plc[index * 4 : index * 4 + 4], "little")
            for index in range(piece_count + 1)
        ]
        pcd_base = (piece_count + 1) * 4
        parts: list[str] = []
        for index in range(piece_count):
            pcd = plc[pcd_base + index * 8 : pcd_base + index * 8 + 8]
            fc = int.from_bytes(pcd[2:6], "little")
            compressed = bool(fc & 0x40000000)
            offset = (fc & 0x3FFFFFFF) // 2 if compressed else (fc & 0x3FFFFFFF)
            char_count = cps[index + 1] - cps[index]
            byte_count = char_count * (1 if compressed else 2)
            raw = word_stream[offset : offset + byte_count]
            parts.append(raw.decode("gb18030" if compressed else "utf-16-le", errors="replace"))
        return _clean_ole_text("".join(parts))


def _scan_unicode_text_runs(data: bytes) -> list[str]:
    """从二进制流中找出连续可读的 UTF-16LE 文本片段。"""
    runs: list[str] = []
    current: list[str] = []
    for index in range(0, len(data) - 1, 2):
        code = int.from_bytes(data[index : index + 2], "little")
        char = chr(code)
        printable = char.isprintable() or char in "\r\n\t"
        if printable and not (0xD800 <= code <= 0xDFFF):
            current.append(char)
        else:
            if len(current) >= 2:
                runs.append("".join(current))
            current = []
    if len(current) >= 2:
        runs.append("".join(current))
    return runs


def _extract_ole_ppt_text(path: Path) -> str:
    """不依赖 Office，从 OLE 演示文稿流中扫描文本（仅作兜底，正文以 COM 为准）。"""
    import olefile

    with olefile.OleFileIO(str(path)) as ole:
        if not ole.exists("PowerPoint Document"):
            raise ValueError("OLE 文件中没有 PowerPoint Document 流。")
        data = ole.openstream("PowerPoint Document").read()

    lines: list[str] = []
    for run in _scan_unicode_text_runs(data):
        text = run.strip()
        if len(text) < 4:
            continue
        cjk_count = sum(1 for ch in text if "\u4e00" <= ch <= "\u9fff")
        if cjk_count < 1 and not any(ch.isalpha() for ch in text):
            continue
        lines.append(text)
    return "\n".join(lines)


def _load_ole_office(path: Path, app_progids: tuple[str, ...], temp_name: str) -> list[Document]:
    try:
        temp_path = _com_save_as(path, app_progids, temp_name, save_format=12)
    except RuntimeError as exc:
        logger.warning("Office/WPS COM 打开失败，回退到内置 OLE 解析: %s", path.name)
        try:
            text = _extract_ole_doc_text(path)
        except Exception as fallback_exc:
            raise ValueError(f"{exc}；内置解析失败: {fallback_exc}") from exc
        if not text.strip():
            raise ValueError(f"{exc}；内置解析未能提取到正文。") from exc
        doc = _text_document(path, text)
        doc.metadata["converted_from"] = path.suffix.lower()
        return [doc]
    try:
        docs = _load_docx(temp_path)
        for doc in docs:
            doc.metadata["source"] = path.name
            doc.metadata["converted_from"] = path.suffix.lower()
        return docs
    finally:
        shutil.rmtree(temp_path.parent, ignore_errors=True)


def _load_doc(path: Path) -> list[Document]:
    fmt = _sniff_format(path)
    if fmt == "zip":
        return _load_zip_document(path)
    if fmt in {"rtf", "html", "text"}:
        return _load_plain_office_text(path, fmt)
    return _load_ole_office(path, ("Word.Application", "KWPS.Application"), "converted.docx")


def _ppt_shape_text(shape) -> str:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    parts: list[str] = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            text = _ppt_shape_text(child)
            if text:
                parts.append(text)
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if text:
                parts.append(text)
    if shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            cells = [cell for cell in cells if cell]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _load_pptx(path: Path) -> list[Document]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    docs: list[Document] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            text = _ppt_shape_text(shape)
            if text:
                lines.append(text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"备注：{notes}")
        if lines:
            docs.append(_text_document(path, "\n".join(lines), slide=slide_number))
    if not docs:
        docs.append(_text_document(path, ""))
    return docs


def _load_ppt(path: Path) -> list[Document]:
    fmt = _sniff_format(path)
    if fmt == "zip":
        return _load_zip_document(path)
    if fmt in {"rtf", "html", "text"}:
        return _load_plain_office_text(path, fmt)
    try:
        temp_path = _com_save_as(
            path,
            ("PowerPoint.Application", "KWPP.Application"),
            "converted.pptx",
            save_format=24,
        )
    except RuntimeError as exc:
        logger.warning("Office/WPS COM 打开失败，回退到内置文本扫描: %s", path.name)
        try:
            text = _extract_ole_ppt_text(path)
        except Exception as fallback_exc:
            raise ValueError(f"{exc}；内置解析失败: {fallback_exc}") from exc
        if not text.strip():
            raise ValueError(f"{exc}；内置解析未能提取到正文。") from exc
        doc = _text_document(path, text)
        doc.metadata["converted_from"] = path.suffix.lower()
        return [doc]
    try:
        docs = _load_pptx(temp_path)
        for doc in docs:
            doc.metadata["source"] = path.name
            doc.metadata["converted_from"] = path.suffix.lower()
        return docs
    finally:
        shutil.rmtree(temp_path.parent, ignore_errors=True)


def _load_wps(path: Path) -> list[Document]:
    fmt = _sniff_format(path)
    if fmt == "zip":
        return _load_zip_document(path)
    if fmt in {"rtf", "html", "text"}:
        return _load_plain_office_text(path, fmt)
    return _load_ole_office(path, ("KWPS.Application", "Word.Application"), "converted.docx")


def _format_cell_value(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.isoformat(sep=" ")
    if isinstance(value, date):
        return value.isoformat()
    return str(value).strip()


def _load_xlsx(path: Path) -> list[Document]:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    docs: list[Document] = []
    try:
        for sheet in workbook.worksheets:
            lines = [f"工作表：{sheet.title}"]
            for index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                cells = [_format_cell_value(value) for value in row]
                cells = [cell for cell in cells if cell]
                if cells:
                    lines.append(f"第{index}行：" + " | ".join(cells))
            docs.append(
                _text_document(
                    path,
                    "\n".join(lines),
                    section=f"工作表：{sheet.title}",
                    section_path=sheet.title,
                )
            )
    finally:
        workbook.close()
    return docs


def _xls_cell_text(sheet, workbook, row_index: int, col_index: int) -> str:
    import xlrd

    cell = sheet.cell(row_index, col_index)
    if cell.ctype == xlrd.XL_CELL_DATE:
        try:
            value = xlrd.xldate.xldate_as_datetime(cell.value, workbook.datemode)
            return _format_cell_value(value)
        except Exception:
            pass
    return _format_cell_value(cell.value)


def _load_xls(path: Path) -> list[Document]:
    import xlrd

    workbook = xlrd.open_workbook(str(path))
    docs: list[Document] = []
    for sheet in workbook.sheets():
        lines = [f"工作表：{sheet.name}"]
        for row_index in range(sheet.nrows):
            cells = [
                _xls_cell_text(sheet, workbook, row_index, col_index)
                for col_index in range(sheet.ncols)
            ]
            cells = [cell for cell in cells if cell]
            if cells:
                lines.append(f"第{row_index + 1}行：" + " | ".join(cells))
        docs.append(
            _text_document(
                path,
                "\n".join(lines),
                section=f"工作表：{sheet.name}",
                section_path=sheet.name,
            )
        )
    return docs


_LOADERS = {
    ".pdf": _load_pdf,
    ".jpg": _load_image,
    ".jpeg": _load_image,
    ".png": _load_image,
    ".txt": _load_text_file,
    ".md": _load_text_file,
    ".csv": _load_csv,
    ".docx": _load_docx,
    ".doc": _load_doc,
    ".ppt": _load_ppt,
    ".pptx": _load_pptx,
    ".wps": _load_wps,
    ".xlsx": _load_xlsx,
    ".xlsm": _load_xlsx,
    ".xls": _load_xls,
}


def load_file(path: Path | str) -> list[Document]:
    """读取单个文件，返回一个或多个 Document（Excel 每个工作表一个）。"""
    path = Path(path)
    suffix = path.suffix.lower()
    loader = _LOADERS.get(suffix)
    if loader is None:
        raise ValueError(f"不支持的文件类型: {path.name}")
    docs = loader(path)
    tier = _tier_from_path(path)
    if tier:
        for doc in docs:
            doc.metadata["tier"] = tier
    return docs
